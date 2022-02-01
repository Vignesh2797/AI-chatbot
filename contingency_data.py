import pptx
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import os
import re
from database import getConn
from itertools import chain

fname = "./static/assets/configuration/Mainline CPT (8th May 21).pptm"
n = 0
Final_image_index=[]
byte_list = []

def update_image_data():
    iter_picture_shapes(Presentation(fname))
    for i, val in enumerate(Master_sub_id):
        Master_sub_id[i] = float(val)
    conn = getConn()
    cur = conn.cursor()
    cur.execute("SET SEARCH_PATH to chatbot")
    for i, val in enumerate(Master_sub_id):
        id = int(val)
        cur.execute(
            "UPDATE ContingencyPlan SET blockage_image = %s WHERE id=%s AND subplan=%s",
            (byte_list[i], id, val),
        )
        conn.commit()
    conn.close()

def write_image(shape):
    global n
    global finder
    slide_index = float(Final_image_index[n])
    image = shape.image
    image_bytes = image.blob
    image_byte_list.append(image_bytes)
    

    image_filename = "image" + str(slide_index) + "." + image.ext
    n += 1


    with open("./static/assets/plans/" + image_filename, "wb") as f:
        f.write(image_bytes)
    f.close()

    to_remove = "./static/assets/plans/" + image_filename
    img_file_size = os.stat("./static/assets/plans/" + image_filename).st_size
    if img_file_size == 14634 or img_file_size==823:
        image_byte_list.remove(image_byte_list[len(image_byte_list)-1])
        os.remove(to_remove)
        n -= 1

    return image_byte_list

def visitor(shape):
    global byte_list
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:

            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        byte_list = write_image(shape)

def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                visitor(shape)
            except:
                continue




def text_extractor(fname, slide_num):
    prs = pptx.Presentation(fname)
    text_runs = []
    for i, sld in enumerate(prs.slides, start=1):
        if i == slide_num:
            for shp in sld.shapes:
                if shp.has_text_frame:
                    print(f"-- slide number {i} --")
                if shp.has_table:
                    tbl = shp.table
                    row_count = len(tbl.rows)
                    col_count = len(tbl.columns)
                    for r in range(0, row_count):
                        text = ""
                        for c in range(0, col_count):
                            cell = tbl.cell(r, c)
                            paragraphs = cell.text_frame.paragraphs
                            for paragraph in paragraphs:
                                for run in paragraph.runs:
                                    text_runs.append(run.text)

    return text_runs


sub_lists = [
    "Status: ",
    "Principles of Service Alteration",
    "Alternative Passenger Journey",
    "Information for Signallers",
    "Additional Information for Station Staff",
    "Passenger Information",
]


def return_index(sub_lists):
    index_list = []
    for i in sub_lists:
        for j in result_ext:
            if i == j:
                index = result_ext.index(i)
                index_list.append(index)
    return index_list


def store_data():
    conn = getConn()
    cur = conn.cursor()
    cur.execute("SET SEARCH_PATH to chatbot")
    
    for i, val in enumerate(Master_sub_id):
        val = float(val)
        c_plan = int(val)
        cur.execute(
            "INSERT INTO ServiceAlteration (cplan_id, subplan, section, alteration) VALUES (%s, %s, %s, %s)",
            [c_plan, val, "Principle", principle_list[i]],
        )
        conn.commit()

   
    for i, val in enumerate(Master_sub_id):
        val = float(val)
        c_plan = int(val)
        cur.execute(
            "INSERT INTO AltPassengerJourney (cplan_id, subplan, section, altjourney) VALUES (%s, %s, %s, %s)",
            [c_plan, val, "Alternate", Alternative_list[i]],
        )
        conn.commit()

    
    for i, val in enumerate(Master_sub_id):
        val = float(val)
        c_plan = int(val)
        cur.execute(
            "INSERT INTO SignallerInfo (cplan_id, subplan, section, info) VALUES (%s, %s, %s, %s)",
            [c_plan, val, "signaller", Information_list[i]],
        )
        conn.commit()

   
    for i, val in enumerate(Master_sub_id):
        val = float(val)
        c_plan = int(val)
        cur.execute(
            "INSERT INTO StationStaffInfo (cplan_id, subplan, section, info) VALUES (%s, %s, %s, %s)",
            [c_plan, val, "additionalinfo", Additional_info_list[i]],
        )
        conn.commit()

   
    for i, val in enumerate(Master_sub_id):
        val = float(val)
        c_plan = int(val)
        cur.execute(
            "INSERT INTO PassengerInfo (cplan_id, subplan, section, info) VALUES (%s, %s, %s, %s)",
            [c_plan, val, "passenger", passenger_info_list[i]],
        )
        conn.commit()

    conn.close()


def return_subplanindex(sub_planlist):
    subplan_status = []
    subplanindex_list = []
    for k in sub_planlist:
        subplan_index = result_ext.index(k)
        subplanindex_list.append(subplan_index)
    m = 0
    while m < len(subplanindex_list):
        try:
            val_subplan = result_ext[subplanindex_list[m] : subplanindex_list[m + 1]]
            val_subplan = "".join(val_subplan)
            try:
                val_subplan = "".join(val_subplan)
            except:
                continue
            try:
                val_subplan = val_subplan.replace("✔", "")
            except:
                continue
        except:
            val_subplan = result_ext[subplanindex_list[m] :]
            val_subplan = "".join(val_subplan)
            try:
                val_subplan = val_subplan.replace("✔", "")
            except:
                continue
        m += 1
        subplan_status.append(val_subplan)
    return subplan_status


Master_sub_id = []
status_list = []
principle_list = []
Alternative_list = []
Additional_info_list = []
Information_list = []
passenger_info_list = []
image_byte_list = []


for x in range(2, 64):
    result_ext = text_extractor(fname, x)  
    if result_ext[0] == "Plan":

       
        plan_id = result_ext[2:3]
        plan_id = "".join(plan_id)

       
        if result_ext[3] == "Surbiton" or result_ext[3] == "Hampton":
            from_to = result_ext[3:5]
            from_to = "".join(from_to)
            from_to_split = from_to.split("-")
            first_stat = from_to_split[0]
            last_stat = from_to_split[1]
            Mainline = result_ext[5:6]
            Mainline = "".join(Mainline)
        else:
            from_to = result_ext[3:4]
            from_to = "".join(from_to)
            from_to_split = from_to.split("-")
            first_stat = from_to_split[0]
            last_stat = from_to_split[1]
            Mainline = result_ext[4:5]
            Mainline = "".join(Mainline)

        
        sub_id = []

        try:
            for x in range(1, 10):
                index = result_ext.index(plan_id + "." + str(x))
                return_subplanindex(sub_id)
                sub_id.append(plan_id + "." + str(x))
        except:
            print("part 1 out of index")

        
        for j in sub_id:
            result_ext

        for x in sub_id:
            Master_sub_id.append(x)

       
        Final_subplan_index = return_subplanindex(sub_id)
        for i, val in enumerate(Final_subplan_index):
            Final_subplan_index[i] = "".join(re.sub(r"[^A-Za-z &]+", "", val))

        conn = getConn()
        cur = conn.cursor()
        cur.execute("SET SEARCH_PATH to chatbot")

        
        for i, val in enumerate(sub_id):
            cur.execute(
                "INSERT INTO ContingencyPlan (id, subplan, name_of_line, first_station, last_station, blockage_type) VALUES (%s, %s, %s, %s, %s, %s)",
                [
                    int(plan_id),
                    float(val),
                    Mainline,
                    first_stat,
                    last_stat,
                    Final_subplan_index[i],
                ],
            )
        conn.commit()
        conn.close()

    else:

        val_index = return_index(sub_lists)
        val_index

        
        if len(val_index) == 7: 
            status = result_ext[val_index[1] + 1 : val_index[2]]
            status = "".join(status)
            status = status.replace("or Status:", ",")
        else:
            status = result_ext[val_index[0] + 1 : val_index[1]]
            status = "".join(status)

        
        if len(val_index) == 7:
            Information = result_ext[val_index[4] + 1 : val_index[5]]
        else:
            Information = result_ext[val_index[3] + 1 : val_index[4]]
        

        
        if len(val_index) == 7:
            principle = result_ext[val_index[2] + 1 : val_index[3]]
        else:
            principle = result_ext[val_index[1] + 1 : val_index[2]]

        if len(val_index) == 7:
            Alternative = result_ext[val_index[3] + 1 : val_index[4]]
        else:
            Alternative = result_ext[val_index[2] + 1 : val_index[3]]

      
        if len(val_index) == 7:
            Additional_info = result_ext[val_index[5] : val_index[6]]
        else:
            try:
                Additional_info = result_ext[val_index[4] : val_index[5]]
            except:
                Additional_info = result_ext[val_index[4] :]
        
        Additional_info = str(Additional_info)
        Additional_info = Additional_info.replace("Additional Information for Station Staff", "")

       
        if len(val_index) == 7:
            passenger_info = result_ext[val_index[6] :]
        else:
            try:
                passenger_info = result_ext[val_index[5] :]
            except:
                passenger_info = result_ext[val_index[4] :]
        
        passenger_info = str(passenger_info)
        passenger_info = passenger_info.replace("Passenger Information", "")

        
        principle_list.append(principle)
        Alternative_list.append(Alternative)
        Additional_info_list.append(Additional_info)
        passenger_info_list.append(passenger_info)
        Information_list.append(Information)
        print("Master", Master_sub_id)

store_data()

for j in Master_sub_id:
        Final_image_index.append(j)
update_image_data()




